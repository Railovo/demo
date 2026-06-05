import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, 'outputs')
prs_path = os.path.join(OUT, 'multi_factor_presentation.pptx')
images = {
    'ic_pvalues':'stat_ic_pvalues.png',
    'robustness':'robustness_heatmap_cost_600bps.png',
    'nav':'highres_multi_factor_cum.png',
    'quintile':'highres_filtered_quintile_cum_returns.png'
}

prs = Presentation()
# title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'HS300 Multi-factor Study — Summary'
subtitle.text = 'Reproducible factors, IC, robustness, and liquidity-aware backtest'

# slide: IC & robustness
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title if shapes.title else None
if title_shape:
    title_shape.text = 'IC significance & Robustness'
left = Inches(0.2)
top = Inches(1.0)
try:
    slide.shapes.add_picture(os.path.join(OUT, images['ic_pvalues']), left, top, width=Inches(4.6))
except Exception:
    pass
try:
    slide.shapes.add_picture(os.path.join(OUT, images['robustness']), Inches(5.0), top, width=Inches(4.0))
except Exception:
    pass

# slide: NAV & quintiles
slide = prs.slides.add_slide(slide_layout)
if slide.shapes.title:
    slide.shapes.title.text = 'Portfolio NAV & Quintile Returns'
try:
    slide.shapes.add_picture(os.path.join(OUT, images['nav']), Inches(0.2), Inches(1.0), width=Inches(4.6))
except Exception:
    pass
try:
    slide.shapes.add_picture(os.path.join(OUT, images['quintile']), Inches(5.0), Inches(1.0), width=Inches(4.0))
except Exception:
    pass

# slide: Methodology & Liquidity
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Methodology & Liquidity Assumptions'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Factors: m120 (momentum), m60, vol60, size. Cross-sectional z-score; industry-neutral; monthly rebalance.'
p = body.add_paragraph()
p.text = 'Liquidity: ADV 20-day average; max 20% ADV per rebalance; linear impact cost applied.'
p.level = 1
p = body.add_paragraph()
p.text = 'Robustness: rolling-window tests and parameter grid; results saved in outputs/robustness_param_summary.csv.'
p.level = 1

prs.save(prs_path)
print('Saved', prs_path)
